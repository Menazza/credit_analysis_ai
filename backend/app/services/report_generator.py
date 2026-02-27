"""
Report generation — Word (credit memo), Excel (financial model), PowerPoint (committee deck).
Production-ready styling, institutional format.
"""
from datetime import date
import csv
import json
from io import BytesIO, StringIO
from typing import Any
import zipfile
from docx import Document as DocxDocument
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, Border, Side, PatternFill
from openpyxl.utils import get_column_letter


# ——— Memo sections (no duplicates; liquidity_leverage removed) ———
MEMO_SECTIONS = [
    "executive_summary",
    "transaction_overview",
    "business_description",
    "industry_overview",
    "competitive_position",
    "financial_performance",
    "financial_risk",
    "cash_flow_liquidity",
    "balance_sheet_leverage",
    "stress_testing_results",
    "accounting_disclosure_quality",
    "key_notes_accounting",
    "key_risks",
    "covenants_headroom",
    "security_collateral",
    "internal_rating_rationale",
    "credit_risk_quantification",
    "recommendation_conditions",
    "monitoring_plan",
    "appendices",
]

# Human-readable section titles
MEMO_SECTION_TITLES = {
    "executive_summary": "Executive Summary",
    "transaction_overview": "Transaction Overview",
    "business_description": "Business Risk Assessment",
    "industry_overview": "Industry Overview",
    "competitive_position": "Competitive Position",
    "financial_performance": "Financial Performance",
    "financial_risk": "Financial Risk",
    "cash_flow_liquidity": "Cash Flow & Liquidity",
    "balance_sheet_leverage": "Balance Sheet & Leverage",
    "stress_testing_results": "Stress Testing Results",
    "accounting_disclosure_quality": "Accounting & Disclosure Quality",
    "key_notes_accounting": "Key Notes (Accounting)",
    "key_risks": "Key Risks",
    "covenants_headroom": "Covenants & Headroom",
    "security_collateral": "Security & Collateral",
    "internal_rating_rationale": "Internal Rating Rationale",
    "credit_risk_quantification": "Credit Risk Quantification (PD/LGD/EAD/ECL)",
    "recommendation_conditions": "Recommendation & Conditions",
    "monitoring_plan": "Monitoring Plan",
    "appendices": "Appendices",
}


def _fmt_num(v: Any) -> str:
    if v is None:
        return "N/A"
    if isinstance(v, (int, float)):
        return f"{v:,.2f}" if abs(v) < 1_000_000 else f"{v/1_000_000:,.2f}m"
    return str(v)


def _fmt_pct(v: Any) -> str:
    if not isinstance(v, (int, float)):
        return "N/A"
    return f"{v * 100:.2f}%"


def _add_formatted_paragraphs(doc: "DocxDocument", text: str, style: str | None = None) -> None:
    """Add paragraphs with institutional formatting: bold subheadings, bullets, proper spacing."""
    para_style = style or "Normal"
    bullet_style = "List Bullet" if "List Bullet" in [s.name for s in doc.styles] else para_style
    subheading_patterns = ("key metrics:", "risk flags:", "source notes:", "governance rules:")
    blocks = text.split("\n\n")
    for block in blocks:
        block = block.strip()
        if not block:
            continue
        lines = [ln.strip() for ln in block.split("\n") if ln.strip()]
        bullet_lines = [ln for ln in lines if ln.startswith(("- ", "• ", "* "))]
        first_lower = lines[0].lower() if lines else ""
        is_subheading_first = first_lower.endswith(":") and any(p in first_lower for p in subheading_patterns)
        if is_subheading_first and len(lines) > 1 and bullet_lines:
            # Bold subheading then bullets
            p = doc.add_paragraph()
            r = p.add_run(lines[0])
            r.bold = True
            r.font.size = Pt(10)
            p.paragraph_format.space_before = Pt(8)
            p.paragraph_format.space_after = Pt(4)
            for line in bullet_lines:
                s = line[2:].lstrip() if len(line) > 2 else line
                if s:
                    bp = doc.add_paragraph(s, style=bullet_style)
                    bp.paragraph_format.space_after = Pt(2)
                    bp.paragraph_format.left_indent = Pt(18)
            if doc.paragraphs:
                doc.paragraphs[-1].paragraph_format.space_after = Pt(6)
        elif len(lines) > 1 and bullet_lines:
            for line in lines:
                if line.startswith(("- ", "• ", "* ")):
                    s = line[2:].lstrip()
                else:
                    s = line
                if s:
                    p = doc.add_paragraph(s, style=bullet_style)
                    p.paragraph_format.space_after = Pt(2)
            if doc.paragraphs:
                doc.paragraphs[-1].paragraph_format.space_after = Pt(8)
        else:
            p = doc.add_paragraph(block, style=para_style)
            p.paragraph_format.space_after = Pt(8)
            p.paragraph_format.line_spacing = 1.15


def build_memo_docx(
    company_name: str,
    review_period_end: date | None,
    version_id: str,
    section_texts: dict[str, str],
    rating_grade: str | None = None,
    recommendation: str = "Maintain",
    key_metrics: dict | None = None,
    metric_by_period: dict | None = None,
    facts_by_period: dict | None = None,
) -> BytesIO:
    doc = DocxDocument()
    try:
        for s in [doc.styles["Normal"], doc.styles["Heading 1"], doc.styles["Heading 2"]]:
            s.font.name = "Calibri"
            s.font.size = Pt(11) if s.name == "Normal" else Pt(14 if s.name == "Heading 1" else 12)
    except KeyError:
        pass
    # Title
    title = doc.add_heading("Credit Review Memo", 0)
    title.alignment = WD_ALIGN_PARAGRAPH.LEFT
    title.paragraph_format.space_after = Pt(6)
    # Metadata block
    meta = doc.add_paragraph()
    meta.add_run("Company: ").bold = True
    meta.add_run(f"{company_name}  ")
    meta.add_run("Review period end: ").bold = True
    meta.add_run(f"{review_period_end or 'N/A'}  ")
    meta.add_run("Version: ").bold = True
    meta.add_run(version_id[:8] + "..." if len(version_id) > 12 else version_id)
    meta.paragraph_format.space_after = Pt(4)
    rating_rec = doc.add_paragraph()
    rating_rec.add_run("Internal rating: ").bold = True
    rating_rec.add_run(f"{rating_grade or 'N/A'}  ")
    rating_rec.add_run("Recommendation: ").bold = True
    rating_rec.add_run(recommendation)
    rating_rec.paragraph_format.space_after = Pt(18)
    # Track 5A: Key metrics table + trend table
    if key_metrics or metric_by_period:
        _add_key_metrics_and_trend_tables(doc, key_metrics or {}, metric_by_period or {}, facts_by_period or {})

    # Sections
    for section_key in MEMO_SECTIONS:
        title_text = MEMO_SECTION_TITLES.get(section_key, section_key.replace("_", " ").title())
        h = doc.add_heading(title_text, level=1)
        h.paragraph_format.space_before = Pt(16)
        h.paragraph_format.space_after = Pt(8)
        text = section_texts.get(section_key) or ""
        if not text.strip():
            text = "Content to be completed."
        _add_formatted_paragraphs(doc, text)
    buf = BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf


def _add_key_metrics_and_trend_tables(doc, key_metrics: dict, metric_by_period: dict, facts_by_period: dict) -> None:
    """Track 5A: Insert key metrics table and trend table into memo."""
    periods = sorted(metric_by_period.keys(), reverse=True)[:5] if metric_by_period else []
    if not periods and not key_metrics:
        return
    h = doc.add_heading("Key Metrics", level=2)
    h.paragraph_format.space_before = Pt(12)
    km_keys = ["revenue", "ebitda", "ebitda_margin", "net_debt_to_ebitda", "net_debt_to_ebitda_incl_leases", "interest_cover", "ebitda_to_interest", "current_ratio"]
    km_labels = {"revenue": "Revenue", "ebitda": "EBITDA", "ebitda_margin": "EBITDA %", "net_debt_to_ebitda": "ND/EBITDA", "net_debt_to_ebitda_incl_leases": "ND/EBITDA (incl)", "interest_cover": "Interest Cover", "ebitda_to_interest": "Interest Cover", "current_ratio": "Current Ratio"}
    tbl = doc.add_table(rows=len(km_keys) + 1, cols=len(periods) + 2)
    tbl.style = "Table Grid"
    row0 = tbl.rows[0].cells
    row0[0].text = "Metric"
    for c, pe in enumerate(periods, 1):
        row0[c].text = str(pe)[:7] if hasattr(pe, "isoformat") else str(pe)[:7]
    for r, k in enumerate(km_keys, 1):
        cells = tbl.rows[r].cells
        cells[0].text = km_labels.get(k, k.replace("_", " ").title())
        for c, pe in enumerate(periods, 1):
            v = metric_by_period.get(pe, {}).get(k) or key_metrics.get(k)
            if v is not None:
                cells[c].text = f"{v:,.2f}" if isinstance(v, (int, float)) and abs(v) < 100 else f"{v:,.1f}"
    doc.add_paragraph()
    # Trend: Revenue, EBITDA, ND/EBITDA
    h2 = doc.add_heading("3-Year Trend", level=2)
    h2.paragraph_format.space_before = Pt(8)
    nd_key = "net_debt_to_ebitda_incl_leases" if any(metric_by_period.get(p, {}).get("net_debt_to_ebitda_incl_leases") for p in periods) else "net_debt_to_ebitda"
    trend_keys = [("revenue", "Revenue"), ("ebitda", "EBITDA"), (nd_key, "ND/EBITDA")]
    tbl2 = doc.add_table(rows=4, cols=len(periods) + 2)
    tbl2.style = "Table Grid"
    tbl2.rows[0].cells[0].text = "Line"
    for c, pe in enumerate(periods[:5], 1):
        tbl2.rows[0].cells[c].text = str(pe)[:7] if hasattr(pe, "isoformat") else str(pe)[:7]
    for r, (k, label) in enumerate(trend_keys, 1):
        tbl2.rows[r].cells[0].text = label
        for c, pe in enumerate(periods[:5], 1):
            v = metric_by_period.get(pe, {}).get(k) or facts_by_period.get(pe, {}).get(k) or (facts_by_period.get(pe, {}).get("operating_profit") if k == "ebitda" else None)
            if v is not None and isinstance(v, (int, float)):
                tbl2.rows[r].cells[c].text = f"{v:,.2f}" if abs(v) < 100 else f"{v:,.0f}"
    doc.add_paragraph()


def _period_key(pe) -> str:
    """Normalize period to lookup key (isoformat or str)."""
    return pe.isoformat() if hasattr(pe, "isoformat") else str(pe)


def _extract_note_refs(text: str, limit: int = 8) -> list[str]:
    if not text:
        return []
    import re
    refs = re.findall(r"Note\s+\d+(?:\.\d+)?", text, flags=re.IGNORECASE)
    out: list[str] = []
    for r in refs:
        rr = r.replace("note", "Note")
        if rr not in out:
            out.append(rr)
        if len(out) >= limit:
            break
    return out


def _simple_pdf_from_lines(title: str, lines: list[str]) -> BytesIO:
    """
    Minimal single-page text PDF generator without extra dependencies.
    """
    safe_lines = [l.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)") for l in lines[:44]]
    offsets: list[int] = []
    pdf = BytesIO()
    pdf.write(b"%PDF-1.4\n")

    def _add_obj(data: bytes) -> int:
        offsets.append(pdf.tell())
        obj_no = len(offsets)
        pdf.write(f"{obj_no} 0 obj\n".encode("ascii"))
        pdf.write(data)
        pdf.write(b"\nendobj\n")
        return obj_no

    font_obj = _add_obj(b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>")
    stream_lines = ["BT /F1 11 Tf 50 790 Td", f"({title}) Tj", "0 -18 Td"]
    for ln in safe_lines:
        stream_lines.append(f"({ln[:140]}) Tj")
        stream_lines.append("0 -15 Td")
    stream_lines.append("ET")
    stream = "\n".join(stream_lines).encode("utf-8")
    content_obj = _add_obj(b"<< /Length " + str(len(stream)).encode("ascii") + b" >>\nstream\n" + stream + b"\nendstream")
    pages_ref = _add_obj(b"<< /Type /Pages /Count 1 /Kids [4 0 R] >>")
    _add_obj(
        f"<< /Type /Page /Parent {pages_ref} 0 R /MediaBox [0 0 612 792] /Resources << /Font << /F1 {font_obj} 0 R >> >> /Contents {content_obj} 0 R >>".encode(
            "ascii"
        )
    )
    catalog_ref = _add_obj(f"<< /Type /Catalog /Pages {pages_ref} 0 R >>".encode("ascii"))

    xref_pos = pdf.tell()
    pdf.write(f"xref\n0 {len(offsets)+1}\n".encode("ascii"))
    pdf.write(b"0000000000 65535 f \n")
    for off in offsets:
        pdf.write(f"{off:010d} 00000 n \n".encode("ascii"))
    pdf.write(
        f"trailer\n<< /Size {len(offsets)+1} /Root {catalog_ref} 0 R >>\nstartxref\n{xref_pos}\n%%EOF".encode("ascii")
    )
    pdf.seek(0)
    return pdf


def _categorize_fact_key(key: str) -> str:
    k = (key or "").lower()
    if any(x in k for x in ("revenue", "cost_of", "gross_profit", "operating", "profit", "tax", "comprehensive")):
        return "pnl"
    if any(x in k for x in ("assets", "equity", "liabilities", "cash", "debt", "inventor", "receivable", "payable")):
        return "bs"
    if any(x in k for x in ("cfo", "capex", "investing", "financing", "cash_flow")):
        return "cf"
    return "other"

def _write_formatted_sheet(ws, title: str, rows: list[dict], period_ends: list, period_key_fn) -> None:
    """Write a sheet with header styling and number formats."""
    header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")
    total_fill = PatternFill(start_color="D6DCE4", end_color="D6DCE4", fill_type="solid")
    bold_font = Font(bold=True)
    ws.cell(1, 1, title).font = Font(bold=True, size=14)
    ws.cell(2, 1, "Line Item").font = header_font
    for c, pe in enumerate(period_ends, start=2):
        cell = ws.cell(2, c, period_key_fn(pe))
        cell.font = header_font
        cell.fill = header_fill
    row = 3
    for r in rows:
        label = r.get("label", r.get("canonical_key", r.get("metric_key", "")))
        vals = r.get("values") or {}
        ws.cell(row, 1, label).font = bold_font if "total" in label.lower() else Font()
        for c, pe in enumerate(period_ends, start=2):
            key = period_key_fn(pe)
            val = vals.get(key)
            cell = ws.cell(row, c, val if val is not None else "")
            if val is not None and isinstance(val, (int, float)):
                cell.number_format = "#,##0" if abs(val) >= 1 else "#,##0.00"
        row += 1
    for col in range(1, len(period_ends) + 2):
        ws.column_dimensions[get_column_letter(col)].width = max(18, min(28, len(title) + 4))

def build_financial_model_xlsx(
    company_name: str,
    period_ends: list[date],
    normalized_rows: list[dict],
    metrics_rows: list[dict],
    version_id: str,
    analysis_output: dict[str, Any] | None = None,
) -> BytesIO:
    wb = Workbook()
    period_key_fn = lambda pe: _period_key(pe)
    pnl_rows = [r for r in normalized_rows if _categorize_fact_key(r.get("canonical_key", "")) == "pnl"]
    bs_rows = [r for r in normalized_rows if _categorize_fact_key(r.get("canonical_key", "")) == "bs"]
    cf_rows = [r for r in normalized_rows if _categorize_fact_key(r.get("canonical_key", "")) == "cf"]
    other_rows = [r for r in normalized_rows if _categorize_fact_key(r.get("canonical_key", "")) == "other"]

    header_fill = PatternFill(start_color="2F5496", end_color="2F5496", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF")

    # Sheet 1: Summary
    ws0 = wb.active
    ws0.title = "Summary"
    ws0.cell(1, 1, "Credit Analysis — Financial Model").font = Font(bold=True, size=16)
    ws0.cell(2, 1, f"Company: {company_name}").font = Font(bold=True)
    ws0.cell(3, 1, f"Version: {version_id}")
    agg = (analysis_output or {}).get("aggregation") or {}
    ws0.cell(5, 1, "Internal Rating:").font = Font(bold=True)
    ws0.cell(5, 2, agg.get("rating_grade", "N/A"))
    ws0.cell(6, 1, "Aggregate Score:").font = Font(bold=True)
    ws0.cell(6, 2, str(agg.get("aggregate_score", "N/A")) + "/100")
    ws0.cell(8, 1, "Key Ratios").font = header_font
    for c, pe in enumerate(period_ends, start=2):
        ws0.cell(8, c, period_key_fn(pe)).font = header_font
        ws0.cell(8, c).fill = header_fill
    row = 9
    for r in metrics_rows[:20]:
        ws0.cell(row, 1, r.get("label", r.get("metric_key", ""))).font = Font(bold=True)
        vals = r.get("values") or {}
        for c, pe in enumerate(period_ends, start=2):
            v = vals.get(period_key_fn(pe))
            cell = ws0.cell(row, c, v if v is not None else "")
            if v is not None and isinstance(v, (int, float)):
                cell.number_format = "#,##0.00" if abs(v) < 1000 else "#,##0"
        row += 1
    ws0.column_dimensions["A"].width = 32
    for c in range(2, len(period_ends) + 2):
        ws0.column_dimensions[get_column_letter(c)].width = 14

    # Sheet 2: P&L
    ws1 = wb.create_sheet("Income Statement")
    all_pnl = pnl_rows if pnl_rows else []
    _write_formatted_sheet(ws1, "Income Statement", all_pnl or [{"label": "No P&L data", "values": {}}], period_ends, period_key_fn)

    # Sheet 3: Balance Sheet
    ws2 = wb.create_sheet("Balance Sheet")
    all_bs = bs_rows if bs_rows else []
    _write_formatted_sheet(ws2, "Balance Sheet", all_bs or [{"label": "No BS data", "values": {}}], period_ends, period_key_fn)

    # Sheet 4: Cash Flow
    ws3 = wb.create_sheet("Cash Flow")
    all_cf = cf_rows if cf_rows else []
    _write_formatted_sheet(ws3, "Cash Flow", all_cf or [{"label": "No CF data", "values": {}}], period_ends, period_key_fn)

    # Sheet 5: Full model (legacy)
    ws4 = wb.create_sheet("Full Model")
    ws4.cell(1, 1, "All Line Items").font = Font(bold=True, size=12)
    row = 3
    ws4.cell(row, 1, "Line Item").font = header_font
    for c, pe in enumerate(period_ends, start=2):
        ws4.cell(row, c, period_key_fn(pe)).font = header_font
        ws4.cell(row, c).fill = header_fill
    row += 1
    for r in normalized_rows:
        ws4.cell(row, 1, r.get("label", r.get("canonical_key", "")))
        for c, pe in enumerate(period_ends, start=2):
            v = (r.get("values") or {}).get(period_key_fn(pe))
            cell = ws4.cell(row, c, v if v is not None else "")
            if v is not None and isinstance(v, (int, float)):
                cell.number_format = "#,##0"
        row += 1
    row += 1
    ws4.cell(row, 1, "Ratios & Metrics").font = Font(bold=True, size=11)
    row += 1
    for r in metrics_rows:
        ws4.cell(row, 1, r.get("label", r.get("metric_key", "")))
        for c, pe in enumerate(period_ends, start=2):
            v = (r.get("values") or {}).get(period_key_fn(pe))
            cell = ws4.cell(row, c, v if v is not None else "")
            if v is not None and isinstance(v, (int, float)):
                cell.number_format = "#,##0.00"
        row += 1

    # Additional institutional sheets (keeps legacy sheets intact)
    ws_raw = wb.create_sheet("Raw Extracted AFS")
    _write_formatted_sheet(ws_raw, "Raw Extracted AFS", normalized_rows or [{"label": "No extracted rows", "values": {}}], period_ends, period_key_fn)

    ws_map = wb.create_sheet("Normalized Mapping")
    ws_map.cell(1, 1, "Canonical Key").font = header_font
    ws_map.cell(1, 2, "Label").font = header_font
    ws_map.cell(1, 3, "Period").font = header_font
    ws_map.cell(1, 4, "Value").font = header_font
    mr = 2
    for r in normalized_rows:
        k = r.get("canonical_key", "")
        lbl = r.get("label", "")
        vals = r.get("values") or {}
        for pe in period_ends:
            pk = period_key_fn(pe)
            ws_map.cell(mr, 1, k)
            ws_map.cell(mr, 2, lbl)
            ws_map.cell(mr, 3, pk)
            ws_map.cell(mr, 4, vals.get(pk))
            mr += 1
    for c, w in zip(["A", "B", "C", "D"], [30, 44, 14, 16]):
        ws_map.column_dimensions[c].width = w

    ws_adj = wb.create_sheet("Adjusted Financials")
    ws_adj.cell(1, 1, "Adjusted Financials (Analyst View)").font = Font(bold=True, size=12)
    adj_rows = []
    for r in normalized_rows:
        ck = (r.get("canonical_key") or "").lower()
        if ck in {
            "revenue",
            "operating_profit",
            "profit_after_tax",
            "total_assets",
            "total_liabilities",
            "total_equity",
            "cash_and_cash_equivalents",
        }:
            adj_rows.append(r)
    _write_formatted_sheet(ws_adj, "Adjusted Financials", adj_rows or [{"label": "No adjusted lines", "values": {}}], period_ends, period_key_fn)

    ws_ratio = wb.create_sheet("Ratio Engine")
    _write_formatted_sheet(ws_ratio, "Ratio Engine", metrics_rows or [{"label": "No ratio data", "values": {}}], period_ends, period_key_fn)

    ws_lease = wb.create_sheet("Lease Adjustments")
    lease_rows = [r for r in metrics_rows if "lease" in (r.get("metric_key") or "").lower()]
    _write_formatted_sheet(ws_lease, "Lease Adjustments", lease_rows or [{"label": "No lease-specific metrics", "values": {}}], period_ends, period_key_fn)

    ws_debt = wb.create_sheet("Debt Schedule")
    debt_rows = [r for r in metrics_rows if any(x in (r.get("metric_key") or "").lower() for x in ("debt", "borrow", "interest_cover", "net_debt"))]
    _write_formatted_sheet(ws_debt, "Debt Schedule", debt_rows or [{"label": "No debt schedule metrics", "values": {}}], period_ends, period_key_fn)

    ws_cov = wb.create_sheet("Covenant Testing")
    ws_cov.cell(1, 1, "Covenant").font = header_font
    ws_cov.cell(1, 2, "Value").font = header_font
    ws_cov.cell(1, 3, "Status").font = header_font
    cov = (((analysis_output or {}).get("section_blocks") or {}).get("covenants") or {}).get("key_metrics") or {}
    cov_items = [
        ("Covenant Leverage Max", cov.get("covenant_leverage_max")),
        ("Current ND/EBITDA", cov.get("current_leverage")),
        ("Leverage Breach", cov.get("leverage_breach")),
        ("Covenant Interest Cover Min", cov.get("covenant_interest_cover_min")),
        ("Current Interest Cover", cov.get("current_interest_cover")),
        ("Interest Cover Breach", cov.get("interest_cover_breach")),
        ("Leverage Headroom %", cov.get("leverage_headroom_pct")),
        ("Coverage Headroom %", cov.get("coverage_headroom_pct")),
    ]
    rr = 2
    for name, val in cov_items:
        ws_cov.cell(rr, 1, name)
        ws_cov.cell(rr, 2, val)
        status = "OK"
        if isinstance(val, bool):
            status = "BREACH" if val else "PASS"
        ws_cov.cell(rr, 3, status)
        rr += 1
    ws_cov.column_dimensions["A"].width = 34
    ws_cov.column_dimensions["B"].width = 20
    ws_cov.column_dimensions["C"].width = 14

    ws_stress = wb.create_sheet("Stress Test Engine")
    ws_stress.cell(1, 1, "Scenario").font = header_font
    ws_stress.cell(1, 2, "ND/EBITDA (stressed)").font = header_font
    ws_stress.cell(1, 3, "Interest Cover (stressed)").font = header_font
    ws_stress.cell(1, 4, "Cash After Shock").font = header_font
    scenarios = ((((analysis_output or {}).get("section_blocks") or {}).get("stress") or {}).get("key_metrics") or {}).get("scenarios") or {}
    sr = 2
    for sname, vals in scenarios.items():
        ws_stress.cell(sr, 1, sname)
        ws_stress.cell(sr, 2, vals.get("net_debt_to_ebitda_stressed"))
        ws_stress.cell(sr, 3, vals.get("interest_cover_stressed"))
        ws_stress.cell(sr, 4, vals.get("cash_after_shock"))
        sr += 1
    ws_stress.column_dimensions["A"].width = 30
    ws_stress.column_dimensions["B"].width = 22
    ws_stress.column_dimensions["C"].width = 24
    ws_stress.column_dimensions["D"].width = 18

    ws_score = wb.create_sheet("Rating Scorecard")
    ws_score.cell(1, 1, "Section").font = header_font
    ws_score.cell(1, 2, "Score").font = header_font
    ws_score.cell(1, 3, "Rating").font = header_font
    agg = (analysis_output or {}).get("aggregation") or {}
    breakdown = agg.get("section_breakdown") or {}
    rr = 2
    for sec, d in breakdown.items():
        ws_score.cell(rr, 1, sec.replace("_", " ").title())
        ws_score.cell(rr, 2, d.get("score"))
        ws_score.cell(rr, 3, d.get("rating") or d.get("section_rating"))
        rr += 1
    ws_score.cell(rr + 1, 1, "Final Internal Rating").font = Font(bold=True)
    ws_score.cell(rr + 1, 2, agg.get("rating_grade"))
    ws_score.cell(rr + 2, 1, "Aggregate Score").font = Font(bold=True)
    ws_score.cell(rr + 2, 2, agg.get("aggregate_score"))
    ws_score.column_dimensions["A"].width = 30
    ws_score.column_dimensions["B"].width = 16
    ws_score.column_dimensions["C"].width = 14

    ws_quant = wb.create_sheet("Credit Risk Quant")
    ws_quant.cell(1, 1, "Credit Risk Quantification").font = Font(bold=True, size=12)
    q = (analysis_output or {}).get("credit_risk_quantification") or {}
    quant_rows = [
        ("Internal Rating", q.get("rating_grade")),
        ("IFRS 9 Stage (proxy)", q.get("ifrs9_stage")),
        ("PD", _fmt_pct(q.get("pd"))),
        ("LGD", _fmt_pct(q.get("lgd"))),
        ("Downturn LGD", _fmt_pct(q.get("downturn_lgd"))),
        ("EAD", q.get("ead")),
        ("Drawn Exposure", q.get("drawn_exposure")),
        ("Undrawn Commitments", q.get("undrawn_commitments")),
        ("CCF", q.get("ccf")),
        ("Expected Loss (EL)", q.get("expected_loss")),
        ("Downturn EL", q.get("expected_loss_downturn")),
    ]
    qr = 3
    ws_quant.cell(2, 1, "Metric").font = header_font
    ws_quant.cell(2, 2, "Value").font = header_font
    for k, v in quant_rows:
        ws_quant.cell(qr, 1, k)
        ws_quant.cell(qr, 2, v)
        if isinstance(v, (int, float)):
            ws_quant.cell(qr, 2).number_format = "#,##0.00"
        qr += 1
    ws_quant.column_dimensions["A"].width = 36
    ws_quant.column_dimensions["B"].width = 22

    ws_dash = wb.create_sheet("Graph Dashboard")
    ws_dash.cell(1, 1, "Graph Dashboard Inputs").font = Font(bold=True, size=12)
    ws_dash.cell(3, 1, "Period").font = header_font
    ws_dash.cell(3, 2, "Revenue").font = header_font
    ws_dash.cell(3, 3, "EBITDA").font = header_font
    ws_dash.cell(3, 4, "ND/EBITDA").font = header_font
    dr = 4
    rev_map = {r.get("metric_key"): r.get("values") or {} for r in metrics_rows}
    for pe in period_ends:
        pk = period_key_fn(pe)
        ws_dash.cell(dr, 1, pk)
        ws_dash.cell(dr, 2, next((r.get("values", {}).get(pk) for r in normalized_rows if r.get("canonical_key") == "revenue"), None))
        ws_dash.cell(dr, 3, rev_map.get("ebitda", {}).get(pk))
        ws_dash.cell(dr, 4, rev_map.get("net_debt_to_ebitda_incl_leases", {}).get(pk) or rev_map.get("net_debt_to_ebitda", {}).get(pk))
        dr += 1
    ws_dash.column_dimensions["A"].width = 14
    ws_dash.column_dimensions["B"].width = 16
    ws_dash.column_dimensions["C"].width = 16
    ws_dash.column_dimensions["D"].width = 16

    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def _add_chart_slides(prs, blank_layout, metric_by_period: dict, facts_by_period: dict) -> None:
    """Track 5C: Add deterministic chart slides - Revenue/EBITDA, ND/EBITDA, CFO vs Capex."""
    try:
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
    except ImportError:
        return
    periods = sorted(metric_by_period.keys(), reverse=True)[:5]
    if not periods:
        periods = sorted(facts_by_period.keys(), reverse=True)[:5]
    if not periods:
        return
    period_labels = [str(p)[:7] if hasattr(p, "isoformat") else str(p)[:7] for p in reversed(periods)]

    def _rev(pe): return facts_by_period.get(pe, {}).get("revenue")
    def _ebitda(pe): return metric_by_period.get(pe, {}).get("ebitda") or facts_by_period.get(pe, {}).get("operating_profit")
    def _nd_eb(pe): return metric_by_period.get(pe, {}).get("net_debt_to_ebitda") or metric_by_period.get(pe, {}).get("net_debt_to_ebitda_incl_leases")
    def _cfo(pe): return facts_by_period.get(pe, {}).get("net_cfo")
    def _capex(pe): return abs(facts_by_period.get(pe, {}).get("capex") or 0)

    # Chart 1: Revenue & EBITDA trend
    rev_vals = [_rev(p) for p in reversed(periods)]
    ebitda_vals = [_ebitda(p) for p in reversed(periods)]
    if any(v is not None for v in rev_vals + ebitda_vals):
        chart_data = CategoryChartData()
        chart_data.categories = period_labels
        chart_data.add_series("Revenue (Rm)", [v if v is not None else 0 for v in rev_vals])
        chart_data.add_series("EBITDA (Rm)", [v if v is not None else 0 for v in ebitda_vals])
        s = prs.slides.add_slide(blank_layout)
        tx = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
        tx.text_frame.paragraphs[0].text = "Revenue & EBITDA Trend"
        tx.text_frame.paragraphs[0].font.bold = True
        s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(0.7), Inches(12), Inches(5.8), chart_data)

    # Chart 2: ND/EBITDA trend
    nd_vals = [_nd_eb(p) for p in reversed(periods)]
    if any(v is not None for v in nd_vals):
        chart_data = CategoryChartData()
        chart_data.categories = period_labels
        chart_data.add_series("ND/EBITDA (x)", [v if v is not None else 0 for v in nd_vals])
        s = prs.slides.add_slide(blank_layout)
        tx = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
        tx.text_frame.paragraphs[0].text = "ND/EBITDA Trend"
        tx.text_frame.paragraphs[0].font.bold = True
        s.shapes.add_chart(XL_CHART_TYPE.LINE, Inches(0.5), Inches(0.7), Inches(12), Inches(5.8), chart_data)

    # Chart 3: CFO vs Capex
    cfo_vals = [_cfo(p) for p in reversed(periods)]
    capex_vals = [_capex(p) for p in reversed(periods)]
    if any(v is not None and v != 0 for v in cfo_vals + capex_vals):
        chart_data = CategoryChartData()
        chart_data.categories = period_labels
        chart_data.add_series("Net CFO (Rm)", [v if v is not None else 0 for v in cfo_vals])
        chart_data.add_series("Capex (Rm)", [v if v is not None else 0 for v in capex_vals])
        s = prs.slides.add_slide(blank_layout)
        tx = s.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.5))
        tx.text_frame.paragraphs[0].text = "CFO vs Capex"
        tx.text_frame.paragraphs[0].font.bold = True
        s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(0.7), Inches(12), Inches(5.8), chart_data)


def build_committee_pptx(
    company_name: str,
    rating_grade: str,
    recommendation: str,
    key_drivers: list[str],
    version_id: str,
    section_texts: dict[str, str] | None = None,
    metric_by_period: dict | None = None,
    facts_by_period: dict | None = None,
    credit_risk_quant: dict[str, Any] | None = None,
) -> BytesIO:
    """Build professional committee deck with styled slides, colors, and explanatory content."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    except ImportError:
        buf = BytesIO()
        buf.write(b"PPTX generation requires python-pptx; placeholder.")
        buf.seek(0)
        return buf

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    sections = section_texts or {}
    quant = credit_risk_quant or {}

    DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
    ACCENT = RGBColor(0x2E, 0x75, 0xB6)
    LIGHT_BG = RGBColor(0xF2, 0xF2, 0xF2)
    RED = RGBColor(0xC0, 0x00, 0x00)
    GREEN = RGBColor(0x00, 0x70, 0x40)

    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    def add_title_slide(title_text: str, subtitle: str = ""):
        s = prs.slides.add_slide(blank_layout)
        left = Inches(0.5)
        top = Inches(2.2)
        width = Inches(12.3)
        tx = s.shapes.add_textbox(left, top, width, Inches(1.2))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        if subtitle:
            tx2 = s.shapes.add_textbox(left, Inches(3.6), width, Inches(1))
            tf2 = tx2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle[:400]
            p2.font.size = Pt(18)
            p2.font.color.rgb = ACCENT

    def add_content_slide(title_text: str, body: str, highlight: str | None = None):
        s = prs.slides.add_slide(blank_layout)
        # Title bar
        box = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))  # Rectangle
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BLUE
        box.line.fill.background()
        tx = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.8))
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        # Body
        body_box = s.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12), Inches(5.5))
        tf = body_box.text_frame
        tf.word_wrap = True
        lines = (body or "N/A").split("\n")[:25]
        for i, line in enumerate(lines[:15]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            line = line.strip()
            if len(line) > 120:
                line = line[:117] + "..."
            p.text = line
            p.font.size = Pt(14)
            p.space_after = Pt(8)
        if highlight:
            hl = s.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(12), Inches(0.8))
            tf2 = hl.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = "► " + highlight[:200]
            p2.font.size = Pt(12)
            p2.font.italic = True
            p2.font.color.rgb = ACCENT

    def _truncate_for_slide(t: str, max_len: int = 1200) -> str:
        t = (t or "")[:max_len]
        if len(t) >= max_len:
            t = t[:max_len - 3] + "..."
        return t

    facts = facts_by_period or {}
    metrics = metric_by_period or {}
    periods = sorted((metrics.keys() or facts.keys()), reverse=True)[:5]
    latest = periods[0] if periods else None
    lp = metrics.get(latest, {}) if latest else {}
    lf = facts.get(latest, {}) if latest else {}

    # Slide 1 — Deal Overview
    deal_overview = (
        f"Borrower: {company_name}\n"
        f"Sector: FMCG / Retail\n"
        f"Facility type: Annual Review Facility\n"
        f"Exposure amount: {_fmt_num(lp.get('net_debt_incl_leases') or lp.get('net_debt_ex_leases') or 0)}\n"
        f"Security summary: See Security & Collateral section\n"
        f"Internal rating: {rating_grade}\n"
        f"Recommendation: {recommendation}"
    )
    add_content_slide("Deal Overview", deal_overview)

    # Slide 2 — Executive Summary
    exec_sum = _truncate_for_slide(sections.get("executive_summary", ""), 1200)
    add_content_slide("Executive Summary", exec_sum, "Trajectory and watchpoints for committee decision.")

    # Slide 3 — Business Overview
    business = sections.get("business_description", "") + "\n\n" + sections.get("competitive_position", "")
    add_content_slide("Business Overview", _truncate_for_slide(business, 1300))

    # Slide 4 — Industry Overview
    add_content_slide("Industry Overview", _truncate_for_slide(sections.get("industry_overview", ""), 1300))

    # Slide 5 — Income Statement Performance
    perf = sections.get("financial_performance", "")
    add_content_slide("Income Statement Performance", _truncate_for_slide(perf, 1300))

    # Slide 6 — Cash Flow Analysis
    add_content_slide("Cash Flow Analysis", _truncate_for_slide(sections.get("cash_flow_liquidity", ""), 1300))

    # Slide 7 — Balance Sheet & Leverage
    add_content_slide("Balance Sheet & Leverage", _truncate_for_slide(sections.get("balance_sheet_leverage", ""), 1300))

    # Slide 8 — Liquidity Position
    liquidity_focus = (
        f"Cash on hand: {_fmt_num(lf.get('cash_and_cash_equivalents'))}\n"
        f"Current ratio: {lp.get('current_ratio', 'N/A')}\n"
        f"Quick ratio: {lp.get('quick_ratio', 'N/A')}\n"
        f"Undrawn facilities: {_fmt_num(lp.get('undrawn_facilities') if isinstance(lp.get('undrawn_facilities'), (int, float)) else 0)}\n\n"
        + sections.get("cash_flow_liquidity", "")
    )
    add_content_slide("Liquidity Position", _truncate_for_slide(liquidity_focus, 1300))

    # Slide 9 — Risk Matrix
    risk_lines = []
    for rf in (sections.get("key_risks", "") or "").split("\n"):
        rf = rf.strip()
        if not rf:
            continue
        risk_lines.append(f"- {rf[:95]} | Severity: High | Trend: Watch | Mitigant: Monitoring/Covenants")
        if len(risk_lines) >= 8:
            break
    risk_body = "Risk | Severity | Trend | Mitigant\n" + "\n".join(risk_lines or ["- No explicit risk flags captured"])
    add_content_slide("Risk Matrix", _truncate_for_slide(risk_body, 1300))

    # Slide 10 — Rating Summary
    rating_body = (
        f"Quantitative score: {lp.get('quant_score', 'See scorecard')}\n"
        f"Qualitative overlay: {lp.get('qual_score', 'See scorecard / commentary')}\n"
        f"Final internal rating: {rating_grade}\n"
        f"Probability of default (PD): {_fmt_pct(quant.get('pd'))}\n"
        f"Loss given default (LGD): {_fmt_pct(quant.get('lgd'))}\n"
        f"Exposure at default (EAD): {_fmt_num(quant.get('ead'))}\n"
        f"Expected loss (EL): {_fmt_num(quant.get('expected_loss'))}\n"
        f"Outlook: {'Negative' if (lp.get('net_debt_to_ebitda_incl_leases') or 0) > 5 else 'Stable'}\n\n"
        + _truncate_for_slide(sections.get("internal_rating_rationale", ""), 700)
    )
    add_content_slide("Rating Summary", rating_body)

    # Slide 11 — Recommendation
    rec_text = (
        f"Decision: {recommendation}\n\n"
        f"Conditions precedent:\n{sections.get('recommendation_conditions', 'To be confirmed')}\n\n"
        f"Monitoring requirements:\n{sections.get('monitoring_plan', 'Quarterly review and covenant monitoring')}"
    )
    add_content_slide("Recommendation", _truncate_for_slide(rec_text, 1300))

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def build_rating_output_json(
    rating_grade: str | None,
    pd_band: Any,
    score_breakdown: dict[str, Any] | None,
    overrides: dict[str, Any] | None,
    rationale: dict[str, Any] | None,
    analysis_output: dict[str, Any] | None,
    credit_risk_quant: dict[str, Any] | None = None,
) -> BytesIO:
    agg = (analysis_output or {}).get("aggregation") or {}
    quant = credit_risk_quant or (analysis_output or {}).get("credit_risk_quantification") or {}
    payload = {
        "final_rating": rating_grade or agg.get("rating_grade"),
        "pd_band": pd_band,
        "pd": quant.get("pd"),
        "lgd": quant.get("lgd"),
        "downturn_lgd": quant.get("downturn_lgd"),
        "ead": quant.get("ead"),
        "expected_loss": quant.get("expected_loss"),
        "expected_loss_downturn": quant.get("expected_loss_downturn"),
        "ifrs9_stage": quant.get("ifrs9_stage"),
        "quant_score": agg.get("aggregate_score"),
        "qual_score": agg.get("qualitative_overlay"),
        "overrides": overrides or {},
        "score_breakdown": score_breakdown or {},
        "rationale": rationale or {},
        "watchlist_flag": bool(agg.get("watchlist_flag")) or ("NEGATIVE" in str(agg.get("outlook", "")).upper()),
        "lgd_estimate": agg.get("lgd_estimate"),
        "outlook": agg.get("outlook"),
    }
    buf = BytesIO()
    buf.write(json.dumps(payload, indent=2, default=str).encode("utf-8"))
    buf.seek(0)
    return buf


def build_memo_pdf(
    company_name: str,
    review_period_end: date | None,
    section_texts: dict[str, str],
    rating_grade: str | None = None,
    recommendation: str | None = None,
) -> BytesIO:
    lines = [
        f"Company: {company_name}",
        f"Review period end: {review_period_end or 'N/A'}",
        f"Internal rating: {rating_grade or 'N/A'}",
        f"Recommendation: {recommendation or 'N/A'}",
        "",
    ]
    for key in MEMO_SECTIONS:
        title = MEMO_SECTION_TITLES.get(key, key.replace("_", " ").title())
        lines.append(f"{title}:")
        txt = (section_texts.get(key) or "Content to be completed.").replace("\n", " ")
        lines.append(txt[:1800])
        lines.append("")
    return _simple_pdf_from_lines("Credit Review Memo", lines)


def build_covenant_certificate_txt(company_name: str, analysis_output: dict[str, Any] | None) -> BytesIO:
    cov = (((analysis_output or {}).get("section_blocks") or {}).get("covenants") or {}).get("key_metrics") or {}
    lines = [
        "COVENANT COMPLIANCE CERTIFICATE",
        f"Borrower: {company_name}",
        "",
        f"Leverage covenant max: {cov.get('covenant_leverage_max', 'N/A')}",
        f"Current leverage: {cov.get('current_leverage', 'N/A')}",
        f"Leverage breach: {cov.get('leverage_breach', 'N/A')}",
        f"Interest cover covenant min: {cov.get('covenant_interest_cover_min', 'N/A')}",
        f"Current interest cover: {cov.get('current_interest_cover', 'N/A')}",
        f"Interest cover breach: {cov.get('interest_cover_breach', 'N/A')}",
        f"Leverage headroom %: {cov.get('leverage_headroom_pct', 'N/A')}",
        f"Coverage headroom %: {cov.get('coverage_headroom_pct', 'N/A')}",
    ]
    buf = BytesIO()
    buf.write("\n".join(lines).encode("utf-8"))
    buf.seek(0)
    return buf


def build_cash_flow_stress_xlsx(company_name: str, analysis_output: dict[str, Any] | None) -> BytesIO:
    wb = Workbook()
    ws = wb.active
    ws.title = "Stress Scenarios"
    ws.cell(1, 1, f"Cash Flow Stress Test Model - {company_name}").font = Font(bold=True, size=13)
    headers = ["Scenario", "ND/EBITDA (stressed)", "Interest Cover (stressed)", "ST Debt/Cash (stressed)", "Cash After Shock"]
    for i, h in enumerate(headers, 1):
        ws.cell(3, i, h).font = Font(bold=True)
    scenarios = ((((analysis_output or {}).get("section_blocks") or {}).get("stress") or {}).get("key_metrics") or {}).get("scenarios") or {}
    r = 4
    for name, vals in scenarios.items():
        ws.cell(r, 1, name)
        ws.cell(r, 2, vals.get("net_debt_to_ebitda_stressed"))
        ws.cell(r, 3, vals.get("interest_cover_stressed"))
        ws.cell(r, 4, vals.get("st_debt_to_cash_stressed"))
        ws.cell(r, 5, vals.get("cash_after_shock"))
        r += 1
    for c, w in zip(["A", "B", "C", "D", "E"], [30, 22, 22, 22, 18]):
        ws.column_dimensions[c].width = w
    buf = BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf


def build_risk_dashboard_pdf(company_name: str, section_texts: dict[str, str], rating_grade: str | None) -> BytesIO:
    lines = [
        f"Company: {company_name}",
        f"Internal Rating: {rating_grade or 'N/A'}",
        "",
        "Top Risks (Severity/Trend/Mitigant):",
    ]
    risks = [r.strip() for r in (section_texts.get("key_risks") or "").split("\n") if r.strip()]
    for r in risks[:12]:
        lines.append(f"- {r[:100]} | High | Watch | Covenant + monthly monitoring")
    if not risks:
        lines.append("- Insufficient risk detail available in current extraction.")
    return _simple_pdf_from_lines("Risk Dashboard", lines)


def build_sector_comparison_appendix_txt(company_name: str, section_texts: dict[str, str]) -> BytesIO:
    text = (
        f"Sector Comparison Appendix - {company_name}\n\n"
        "Peer and sector comparator tables should be populated from external market datasets.\n"
        "Current baseline includes internal sector narrative from memo sections.\n\n"
        + (section_texts.get("industry_overview") or "No industry overview available.")
    )
    buf = BytesIO()
    buf.write(text.encode("utf-8"))
    buf.seek(0)
    return buf


def build_data_room_zip(
    company_name: str,
    version_id: str,
    normalized_rows: list[dict],
    metrics_rows: list[dict],
    section_texts: dict[str, str],
    analysis_output: dict[str, Any] | None,
    rating_output: dict[str, Any] | None,
) -> BytesIO:
    """
    Build a structured ZIP bundle for due-diligence / audit handoff.
    Includes core structured outputs and mapping logs.
    """
    buf = BytesIO()
    with zipfile.ZipFile(buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("meta/version.txt", f"company={company_name}\nversion={version_id}\n")
        zf.writestr("structured/analysis_output.json", json.dumps(analysis_output or {}, indent=2, default=str))
        zf.writestr("structured/rating_output.json", json.dumps(rating_output or {}, indent=2, default=str))
        zf.writestr("structured/section_texts.json", json.dumps(section_texts or {}, indent=2, default=str))

        norm_csv = StringIO()
        w = csv.writer(norm_csv)
        w.writerow(["canonical_key", "label", "period", "value"])
        for r in normalized_rows:
            key = r.get("canonical_key", "")
            label = r.get("label", "")
            vals = r.get("values") or {}
            for p, v in vals.items():
                w.writerow([key, label, p, v])
        zf.writestr("mapping_logs/normalized_mapping.csv", norm_csv.getvalue())

        met_csv = StringIO()
        w2 = csv.writer(met_csv)
        w2.writerow(["metric_key", "label", "period", "value"])
        for r in metrics_rows:
            key = r.get("metric_key", "")
            label = r.get("label", "")
            vals = r.get("values") or {}
            for p, v in vals.items():
                w2.writerow([key, label, p, v])
        zf.writestr("mapping_logs/metric_facts.csv", met_csv.getvalue())
    buf.seek(0)
    return buf
